import logging
import os
import threading
from pptx import Presentation
from pptx.util import Inches

MEDIA_FOLDER = "./test"  # Path to the media folder
PPTX_PATH = "t.pptx"  # Path to the PowerPoint presentation
FILE_LOG = "app.log"  # Name of the log file
PRESENTATION_TITLE = "My Media Presentation"  # Title of the presentation

# Set up logging configuration
logging.basicConfig(filename=FILE_LOG, level=logging.INFO)
logging, logging.info(f"file-pptx {PPTX_PATH if PPTX_PATH else 'not pptx name'}"),


# Function to upload media files to the PowerPoint presentation
def upload_media_to_pptx(media_folder, prs):
    print(f"in {FILE_LOG} Record all events")
    try:
        for filename in os.listdir(media_folder):
            if filename.endswith((".mp4", ".jpg", ".png", ".gif", ".wav", ".mp3")):
                media_path = os.path.join(media_folder, filename)
                media_name_with_extension = os.path.basename(media_path)
                # Create a new slide
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                # Set the slide title
                title = slide.shapes.title
                title.text = media_name_with_extension
                try:
                    # Add the media file to the slide
                    if filename.endswith((".mp4", ".wav", ".mp3")):
                        slide.shapes.add_movie(
                            media_path,
                            left=0,
                            top=0,
                            width=prs.slide_width,
                            height=prs.slide_height,
                        )
                    else:
                        pic = slide.shapes.add_picture(
                            media_path,
                            left=Inches(0),
                            top=Inches(0),
                            width=prs.slide_width,
                            height=prs.slide_height,
                        )
                    logging.info(f"Uploaded file: {media_name_with_extension}")
                except Exception as err:
                    logging.error(
                        f"Error uploading file {media_name_with_extension}: {err}"
                    )
            else:
                logging.warning(f"Skipping unsupported file: {filename}")
                continue
    finally:
        msg = f"Saving presentation to {PPTX_PATH}"
        logging.info(msg)
        print(msg)


# Check if the PowerPoint file exists and create/open it
def initialize_presentation(pptx_path):
    if os.path.exists(pptx_path):
        print(f"Adding to an existing PowerPoint file: {pptx_path}")
        prs = Presentation(pptx_path)
    else:
        print(f"Creating new PowerPoint file: {pptx_path}")
        prs = Presentation()
        # Add a title slide with the presentation title
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        title.text = PRESENTATION_TITLE
    return prs


if __name__ == "__main__":
    prs = initialize_presentation(PPTX_PATH)
    try:
        thread = threading.Thread(target=upload_media_to_pptx, args=(MEDIA_FOLDER, prs))
        thread.start()
        thread.join()
        prs.save(PPTX_PATH)
    except Exception as e:
        print(f"Error: {e}")
        logging.warning(f"Error: {e}")
