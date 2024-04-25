import logging
import os
import threading
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor


class UploadMediaPresentation:
    def __init__(
        self, media_folder, pptx_path, log_file="app.log", presentation_title=None
    ):
        self.media_folder = media_folder
        self.pptx_path = pptx_path
        self.presentation_title = (
            presentation_title if presentation_title else f"name pptx\n{pptx_path}"
        )
        self.log_file = log_file
        self.prs = self.initialize_presentation(self.pptx_path)

    def get_pptx_path(self):
        return self.pptx_path

    def get_media_folder(self):
        return self.media_folder

    def get_log_file(self):
        return self.log_file

    def set_log_file(self, log_file):
        self.log_file = log_file

    def clear_log_file(self):
        """
        Clears the content of the log file.
        """
        with open(self.log_file, "w"):
            pass

    def get_log_content(self):
        """
        Reads and returns the content of the log file.
        """
        with open(self.log_file, "r+") as file:
            txt = file.read()
        return txt

    def setup_logging(self):
        """
        Setup logging configuration.
        """
        logging.basicConfig(filename=self.log_file, level=logging.INFO)
        logging.info(
            f"file-pptx {self.pptx_path if self.pptx_path else 'not pptx name'}"
        )

    def upload_media_to_pptx(self):
        """
        Upload media files to the PowerPoint presentation.
        """
        print(f"in {self.log_file} Record all events")
        try:
            for filename in os.listdir(self.media_folder):
                if filename.endswith((".mp4", ".jpg", ".png", ".gif", ".wav", ".mp3")):
                    media_path = os.path.join(self.media_folder, filename)
                    media_name_with_extension = os.path.basename(media_path)
                    slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
                    title = slide.shapes.title
                    title.text = media_name_with_extension
                    try:
                        if filename.endswith((".mp4", ".wav", ".mp3")):
                            slide.shapes.add_movie(
                                media_path,
                                left=0,
                                top=0,
                                width=self.prs.slide_width,
                                height=self.prs.slide_height,
                            )
                        elif filename.endswith((".jpg", ".png", ".gif")):
                            slide.shapes.add_picture(
                                media_path,
                                left=Inches(0),
                                top=Inches(0),
                                width=self.prs.slide_width,
                                height=self.prs.slide_height,
                            )
                        else:
                            raise ValueError("Unsupported image type.")
                        logging.info(f"Uploaded file: {media_name_with_extension}")
                    except Exception as err:
                        logging.error(
                            f"Error uploading file {media_name_with_extension}: {err}"
                        )
                else:
                    logging.warning(f"Skipping unsupported file: {filename}")
                    continue
        finally:
            msg = f"Saving presentation to {self.pptx_path}"
            logging.info(msg)
            print(msg)

    def set_text_color(self, text_frame, color):
        """
        Set the font color of all runs in all paragraphs of a text frame.

        Parameters:
        - text_frame: The text frame whose text color is to be set.
        - color: The RGB color to set (e.g., RGBColor(255, 0, 0) for red).
        """
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = color

    def initialize_presentation(self, pptx_path):
        """
        Initialize a new PowerPoint presentation or open an existing one.

        Parameters:
        - pptx_path: Path to the PowerPoint presentation.

        Returns:
        - prs: PowerPoint presentation object.
        """
        if os.path.exists(pptx_path):
            print(f"Adding to an existing PowerPoint file: {pptx_path}")
            prs = Presentation(pptx_path)
        else:
            print(f"Creating new PowerPoint file: {pptx_path}")
            prs = Presentation()
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = title_slide.shapes.title
            title.text = self.presentation_title
            self.set_text_color(title.text_frame, RGBColor(0, 128, 128))

        return prs

    def start(self):
        """
        Start the process of uploading media to the PowerPoint presentation.
        """
        self.setup_logging()
        try:
            thread = threading.Thread(target=self.upload_media_to_pptx)
            thread.start()
            thread.join()
            self.prs.save(self.pptx_path)
        except Exception as e:
            print(f"Error: {e}")
            logging.warning(f"Error: {e}")


if __name__ == "__main__":
    MEDIA_FOLDER = "./test"
    PPTX_PATH = "t.pptx"

    media_presentation = UploadMediaPresentation(
        media_folder=MEDIA_FOLDER, pptx_path=PPTX_PATH
    )
    media_presentation.start()
