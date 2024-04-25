import logging
import os
import shutil
from pptx import Presentation
from pptx.util import Inches

MEDIA_FOLDER = "./Output_folder/"  # Path to the media folder
PPTX_PATH = "t.pptx"  # Path to the PowerPoint presentation
MEDIA_SAVE_FOLDER = "media_saved"  # Folder to save media files
FILE_LOG = "app3.log"  # Name of the log file
PRESENTATION_TITLE = "My Media Presentation"  # Title of the presentation

# Set up logging configuration
logging.basicConfig(filename=FILE_LOG, level=logging.INFO)
logging.info(f"file-pptx {PPTX_PATH if PPTX_PATH else 'not pptx name'}")


# Function to upload media files to the PowerPoint presentation
def upload_media_to_pptx(media_folder, prs, save_folder):
    print(f"in {FILE_LOG} Record all events")
    try:
        for filename in os.listdir(media_folder):
            if filename.endswith((".mp4", ".jpg", ".png", ".gif", ".wav", ".mp3")):
                media_path = os.path.join(media_folder, filename)
                media_name_with_extension = os.path.basename(media_path)
                # Create a new slide
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                # Set the slide title
                title = slide.shapes.title
                title.text = media_name_with_extension
                try:
                    # Save the media file to the specified folder
                    shutil.copy(
                        media_path, os.path.join(save_folder, media_name_with_extension)
                    )
                    logging.info(f"Saved file: {media_name_with_extension}")
                except Exception as err:
                    logging.error(
                        f"Error saving file {media_name_with_extension}: {err}"
                    )
            else:
                logging.warning(f"Skipping unsupported file: {filename}")
                continue
    finally:
        msg = f"Saving presentation to {PPTX_PATH}"
        logging.info(msg)
        print(msg)


# Check if the PowerPoint file exists and create/open it
def initialize_presentation(pptx_path):
    if os.path.exists(pptx_path):
        print(f"Adding to an existing PowerPoint file: {pptx_path}")
        prs = Presentation(pptx_path)
    else:
        print(f"Creating new PowerPoint file: {pptx_path}")
        prs = Presentation()
        # Add a title slide with the presentation title
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        title.text = PRESENTATION_TITLE
    return prs


if __name__ == "__main__":
    prs = initialize_presentation(PPTX_PATH)
    try:
        # Create the folder to save media files if it doesn't exist
        os.makedirs(MEDIA_SAVE_FOLDER, exist_ok=True)
        upload_media_to_pptx(MEDIA_FOLDER, prs, MEDIA_SAVE_FOLDER)
        prs.save(PPTX_PATH)
    except Exception as e:
        print(f"Error: {e}")
        logging.warning(f"Error: {e}")
