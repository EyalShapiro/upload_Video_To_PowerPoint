from array import array
import logging
import os
import threading
from pptx import Presentation

##const
VIDEO_FOLDER = "test"
PPTX_PATH = "test/my-PowerPoint.pptx"
# -----------------
logging.basicConfig(filename="app.log", level=logging.INFO)

prs = Presentation()


def upload_Video_To_PPTX(video_name):
    print("in app.log Record all events")

    for filename in os.listdir(video_name):
        if filename.endswith(".mp4"):
            video_path = os.path.join(video_name, filename)
            video_name_without_extension = os.path.splitext(filename)[0]
            # יצירת עמוד חדש במצגת
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            # הוספת כותרת לעמוד
            title = slide.shapes.title
            title.text = video_name_without_extension
            try:
                # הוספת הוידאו לעמוד
                slide.shapes.add_movie(
                    video_path,
                    left=0,
                    top=0,
                    width=prs.slide_width,
                    height=prs.slide_height,
                )
                logging.info(f"Uploaded file: {video_name_without_extension}")
                print(f"Uploaded file: {video_name_without_extension}")
            except Exception as e:

                logging.error(
                    f"Error uploading file {video_name_without_extension}: {e}"
                )
                print(f"Error uploading file {video_name_without_extension}: {e}")


def File_Is_Existed(PPTX_PATH):
    if os.path.exists(PPTX_PATH):
        print("Adding to an existing PowerPoint file:{PPTX_PATH}")
    else:
        print(f"cerate new PowerPoint file:\t{PPTX_PATH}")


if __name__ == "__main__":
    File_Is_Existed(PPTX_PATH)
    try:
        thread = threading.Thread(target=upload_Video_To_PPTX, args=(VIDEO_FOLDER,))
        thread.start()
        thread.join()
        prs.save(PPTX_PATH)
    except Exception as e:
        logging.error(f"Error: {e}")
