import logging
import os
import threading
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor


class MediaPresentation:
    def __init__(
        self, media_folder, pptx_path, log_file="app.log", presentation_title=None
    ):
        self.MEDIA_FOLDER = media_folder
        self.PPTX_PATH = pptx_path
        self.presentation_title = (
            presentation_title if presentation_title else "name pptx\n" + pptx_path
        )
        self.log_file = log_file
        self.prs = self.Initialize_Presentation(self.PPTX_PATH)

    # Getter and setter methods for log_file
    def Get_Log_File(self):
        return self.log_file

    def Set_Log_File(self, log_file):
        self.log_file = log_file

    # Method to clear the log file
    def Clear_Log_File(self):
        """
        Clears the content of the log file.
        """
        with open(self.log_file, "w"):
            pass

    # Method to get the log content
    def Get_Log_Content(self):
        """
        Reads and returns the content of the log file.
        """
        with open(self.log_file, "r") as f:
            return f.read()

    def Setup_Logging(self):
        """
        Setup logging configuration.
        """
        logging.basicConfig(filename=self.log_file, level=logging.INFO)
        logging.info(
            f"file-pptx {self.PPTX_PATH if self.PPTX_PATH else 'not pptx name'}"
        )

    def Upload_Media_To_PPTX(self):
        """
        Upload media files to the PowerPoint presentation.
        """
        print(f"in {self.log_file} Record all events")
        try:
            for filename in os.listdir(self.MEDIA_FOLDER):
                if filename.endswith((".mp4", ".jpg", ".png", ".gif", ".wav", ".mp3")):
                    media_path = os.path.join(self.MEDIA_FOLDER, filename)
                    media_name_with_extension = os.path.basename(media_path)
                    # Create a new slide
                    slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
                    # Set the slide title
                    title = slide.shapes.title
                    title.text = media_name_with_extension
                    try:
                        # Add the media file to the slide
                        if filename.endswith((".mp4", ".wav", ".mp3")):
                            slide.shapes.add_movie(
                                media_path,
                                left=0,
                                top=0,
                                width=self.prs.slide_width,
                                height=self.prs.slide_height,
                            )
                        else:
                            pic = slide.shapes.add_picture(
                                media_path,
                                left=Inches(0),
                                top=Inches(0),
                                width=self.prs.slide_width,
                                height=self.prs.slide_height,
                            )
                        logging.info(f"Uploaded file: {media_name_with_extension}")
                    except Exception as err:
                        logging.error(
                            f"Error uploading file {media_name_with_extension}: {err}"
                        )
                else:
                    logging.warning(f"Skipping unsupported file: {filename}")
                    continue
        finally:
            msg = f"Saving presentation to {self.PPTX_PATH}"
            logging.info(msg)
            print(msg)

    def Set_Text_Color(self, text_frame, color):
        """
        Set the font color of all runs in all paragraphs of a text frame.

        Parameters:
        - text_frame: The text frame whose text color is to be set.
        - color: The RGB color to set (e.g., RGBColor(255, 0, 0) for red).
        """
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = color

    def Initialize_Presentation(self, pptx_path):
        """
        Initialize a new PowerPoint presentation or open an existing one.

        Parameters:
        - pptx_path: Path to the PowerPoint presentation.

        Returns:
        - prs: PowerPoint presentation object.
        """
        if os.path.exists(pptx_path):
            print(f"Adding to an existing PowerPoint file: {pptx_path}")
            prs = Presentation(pptx_path)
        else:
            print(f"Creating new PowerPoint file: {pptx_path}")
            prs = Presentation()
            # Add a title slide with the presentation title
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = title_slide.shapes.title
            title.text = self.presentation_title
            self.Set_Text_Color(
                title.text_frame, RGBColor(0, 128, 128)
            )  # Set text color to teal

        return prs

    def Start(self):
        """
        Start the process of uploading media to the PowerPoint presentation.
        """
        self.Setup_Logging()
        try:
            thread = threading.Thread(target=self.Upload_Media_To_PPTX)
            thread.start()
            thread.join()
            self.prs.save(self.PPTX_PATH)
        except Exception as e:
            print(f"Error: {e}")
            logging.warning(f"Error: {e}")


if __name__ == "__main__":
    MEDIA_FOLDER = "./test"  # Path to the media folder
    PPTX_PATH = "t.pptx"  # Path to the PowerPoint presentation

    media_presentation = MediaPresentation(
        media_folder=MEDIA_FOLDER, pptx_path=PPTX_PATH
    )
    media_presentation.Start()
